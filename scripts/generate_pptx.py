"""Generate a professional PowerPoint presentation for the Enterprise Security Guardrail Auditor."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colors ──
NAVY = RGBColor(0x0F, 0x17, 0x2A)
DARK_BLUE = RGBColor(0x1E, 0x29, 0x3B)
ACCENT_BLUE = RGBColor(0x38, 0x82, 0xF6)
ACCENT_TEAL = RGBColor(0x14, 0xB8, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
YELLOW = RGBColor(0xEA, 0xB3, 0x08)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
SUBTLE_TEXT = RGBColor(0x64, 0x74, 0x8B)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_bg(slide, color=NAVY):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    """Add a text box with a single run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return tf


def add_bullet_list(
    slide,
    left,
    top,
    width,
    height,
    items,
    font_size=16,
    color=WHITE,
    bullet_color=ACCENT_BLUE,
    spacing=Pt(8),
):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        # bullet char
        bullet_run = p.add_run()
        bullet_run.text = "●  "
        bullet_run.font.size = Pt(font_size)
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.name = "Calibri"
        # text
        text_run = p.add_run()
        text_run.text = item
        text_run.font.size = Pt(font_size)
        text_run.font.color.rgb = color
        text_run.font.name = "Calibri"
    return tf


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_card(slide, left, top, width, height, title, value, accent=ACCENT_BLUE):
    """Add a metric card with a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    # Title
    add_text_box(
        slide,
        left + Inches(0.2),
        top + Inches(0.15),
        width - Inches(0.4),
        Inches(0.4),
        title,
        font_size=12,
        color=MID_GRAY,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )
    # Value
    add_text_box(
        slide,
        left + Inches(0.2),
        top + Inches(0.5),
        width - Inches(0.4),
        Inches(0.6),
        value,
        font_size=28,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def add_table(
    slide, left, top, width, rows_data, col_widths=None, header_color=ACCENT_BLUE
):
    """Add a styled table."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(
        rows, cols, left, top, width, Inches(0.45 * rows)
    )
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = "Calibri"

            if r == 0:
                # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                run.font.size = Pt(13)
                run.font.color.rgb = WHITE
                run.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE if r % 2 == 1 else NAVY
                run.font.size = Pt(12)
                run.font.color.rgb = LIGHT_GRAY
                p.alignment = PP_ALIGN.LEFT

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Remove cell borders by making them transparent
            from pptx.oxml.ns import qn

            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
                ln = tcPr.find(qn(border_name))
                if ln is None:
                    from lxml import etree

                    ln = etree.SubElement(tcPr, qn(border_name))
                ln.set("w", "0")
                noFill = ln.find(qn("a:noFill"))
                if noFill is None:
                    from lxml import etree

                    etree.SubElement(ln, qn("a:noFill"))


def add_section_title(slide, text, subtitle=None):
    """Standard section title layout."""
    add_bg(slide, NAVY)
    add_text_box(
        slide,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(0.6),
        text,
        font_size=32,
        color=WHITE,
        bold=True,
    )
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_BLUE)
    if subtitle:
        add_text_box(
            slide,
            Inches(0.8),
            Inches(1.4),
            Inches(11),
            Inches(0.5),
            subtitle,
            font_size=16,
            color=MID_GRAY,
        )


# ── Slide Builders ──


def slide_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, NAVY)

    # Top accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # Shield icon placeholder (text emoji)
    add_text_box(
        slide,
        Inches(5.2),
        Inches(1.5),
        Inches(3),
        Inches(1.2),
        "🛡️",
        font_size=64,
        color=ACCENT_BLUE,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Inches(1.5),
        Inches(2.7),
        Inches(10.3),
        Inches(1),
        "Enterprise Security Guardrail Auditor",
        font_size=40,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_accent_line(slide, Inches(5.2), Inches(3.8), Inches(3), ACCENT_BLUE)

    add_text_box(
        slide,
        Inches(2),
        Inches(4.1),
        Inches(9.3),
        Inches(0.6),
        "API-first Infrastructure Security Scanner  •  Terraform  •  Risk Scoring  •  Dashboard",
        font_size=18,
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Inches(2),
        Inches(5.2),
        Inches(9.3),
        Inches(0.5),
        "Technical Architecture Review & MVP Demo",
        font_size=16,
        color=SUBTLE_TEXT,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Inches(2),
        Inches(6.3),
        Inches(9.3),
        Inches(0.5),
        "Adedoyin Simeon Adeyemi  •  May 2026",
        font_size=14,
        color=SUBTLE_TEXT,
        alignment=PP_ALIGN.CENTER,
    )

    # Bottom accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def slide_problem(prs):
    """Slide 2: Problem Statement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide,
        "Problem Statement",
        "Infrastructure misconfigurations are the #1 cause of cloud security breaches",
    )

    problems = [
        ("☁️", "Public S3 Buckets", "Exposing sensitive data to the internet"),
        ("🔓", "Open SSH Ports", "Inviting brute-force attacks on servers"),
        ("💾", "Unencrypted Databases", "Violating compliance (HIPAA, SOC2, PCI-DSS)"),
        ("👤", "Wildcard IAM Policies", "Granting god-mode access across accounts"),
    ]

    for i, (icon, title, desc) in enumerate(problems):
        x = Inches(0.8 + i * 3.05)
        y = Inches(2.3)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(3.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BLUE
        card.line.color.rgb = RGBColor(0x33, 0x44, 0x55)
        card.line.width = Pt(1)

        add_text_box(
            slide,
            x + Inches(0.2),
            y + Inches(0.3),
            Inches(2.4),
            Inches(0.8),
            icon,
            font_size=40,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x + Inches(0.2),
            y + Inches(1.2),
            Inches(2.4),
            Inches(0.5),
            title,
            font_size=18,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x + Inches(0.2),
            y + Inches(1.9),
            Inches(2.4),
            Inches(1.2),
            desc,
            font_size=14,
            color=MID_GRAY,
            alignment=PP_ALIGN.CENTER,
        )

    add_text_box(
        slide,
        Inches(1),
        Inches(6.2),
        Inches(11),
        Inches(0.5),
        "Teams need automated, continuous auditing of infrastructure-as-code before deployment.",
        font_size=16,
        color=ACCENT_TEAL,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def slide_solution(prs):
    """Slide 3: Solution Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide,
        "Solution Overview",
        "Enterprise Security Guardrail Auditor — an API-first infrastructure scanner",
    )

    steps = [
        ("1", "PARSE", "Terraform\nconfiguration files", ACCENT_BLUE),
        ("2", "EVALUATE", "Against 15\nsecurity rules", ACCENT_TEAL),
        ("3", "SCORE", "Weighted risk\n0–100 scale", ORANGE),
        ("4", "VISUALIZE", "Real-time\ndashboard", GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.8 + i * 3.05)
        y = Inches(2.5)

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.9), y, Inches(1), Inches(1)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text_box(
            slide,
            x + Inches(0.9),
            y + Inches(0.15),
            Inches(1),
            Inches(0.7),
            num,
            font_size=32,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Arrow (except last)
        if i < 3:
            arrow_x = x + Inches(2.1)
            add_text_box(
                slide,
                arrow_x,
                y + Inches(0.2),
                Inches(0.8),
                Inches(0.6),
                "→",
                font_size=28,
                color=MID_GRAY,
                alignment=PP_ALIGN.CENTER,
            )

        add_text_box(
            slide,
            x + Inches(0.2),
            y + Inches(1.2),
            Inches(2.4),
            Inches(0.5),
            title,
            font_size=20,
            color=color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x + Inches(0.2),
            y + Inches(1.8),
            Inches(2.4),
            Inches(0.8),
            desc,
            font_size=14,
            color=MID_GRAY,
            alignment=PP_ALIGN.CENTER,
        )

    # Bottom tagline
    add_text_box(
        slide,
        Inches(1),
        Inches(5.5),
        Inches(11),
        Inches(0.8),
        "Upload a .tf file  →  Get instant security findings with remediation guidance",
        font_size=20,
        color=WHITE,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )


def slide_architecture(prs):
    """Slide 4: System Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide, "System Architecture", "API-first, async, pluggable rule engine"
    )

    # Frontend box
    fe_x, fe_y = Inches(0.8), Inches(2.5)
    fe = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, fe_x, fe_y, Inches(4.5), Inches(3.5)
    )
    fe.fill.solid()
    fe.fill.fore_color.rgb = DARK_BLUE
    fe.line.color.rgb = ACCENT_BLUE
    fe.line.width = Pt(2)
    add_text_box(
        slide,
        fe_x + Inches(0.3),
        fe_y + Inches(0.2),
        Inches(3.9),
        Inches(0.5),
        "Frontend — React 19 + TypeScript",
        font_size=16,
        color=ACCENT_BLUE,
        bold=True,
    )
    add_bullet_list(
        slide,
        fe_x + Inches(0.3),
        fe_y + Inches(0.8),
        Inches(3.9),
        Inches(2.5),
        [
            "Dashboard & Risk Visualization",
            "Scan Management",
            "Guardrail Configuration",
            "Tailwind CSS + Recharts",
        ],
        font_size=13,
        color=LIGHT_GRAY,
    )

    # Arrow
    add_text_box(
        slide,
        Inches(5.5),
        Inches(3.8),
        Inches(1.5),
        Inches(0.6),
        "REST API",
        font_size=14,
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(5.5),
        Inches(4.2),
        Inches(1.5),
        Inches(0.5),
        "◀─────▶",
        font_size=18,
        color=ACCENT_BLUE,
        alignment=PP_ALIGN.CENTER,
    )

    # Backend box
    be_x, be_y = Inches(7.3), Inches(2.5)
    be = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, be_x, be_y, Inches(5.2), Inches(3.5)
    )
    be.fill.solid()
    be.fill.fore_color.rgb = DARK_BLUE
    be.line.color.rgb = ACCENT_TEAL
    be.line.width = Pt(2)
    add_text_box(
        slide,
        be_x + Inches(0.3),
        be_y + Inches(0.2),
        Inches(4.6),
        Inches(0.5),
        "Backend — FastAPI + SQLAlchemy Async",
        font_size=16,
        color=ACCENT_TEAL,
        bold=True,
    )
    add_bullet_list(
        slide,
        be_x + Inches(0.3),
        be_y + Inches(0.8),
        Inches(4.6),
        Inches(2.5),
        [
            "Thin Controllers + Service Layer",
            "Scanner: Parser → Rules → Scoring",
            "Async throughout (aiosqlite)",
            "SQLite Database",
        ],
        font_size=13,
        color=LIGHT_GRAY,
    )

    # Key principles
    principles = [
        "API-first: Frontend is a pure REST consumer",
        "Pluggable rules: BaseRule subclasses",
        "Async throughout: FastAPI + SQLAlchemy async",
    ]
    add_bullet_list(
        slide,
        Inches(0.8),
        Inches(6.3),
        Inches(11),
        Inches(1),
        principles,
        font_size=13,
        color=MID_GRAY,
        bullet_color=ACCENT_TEAL,
    )


def slide_scanner(prs):
    """Slide 5: Scanner Engine Design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide, "Scanner Engine Design", "Four-stage pipeline with pluggable rules"
    )

    table_data = [
        ["Stage", "Component", "Output"],
        ["1. Parse", "TerraformParser", "list[ParsedResource]"],
        ["2. Evaluate", "RuleRegistry → BaseRule.evaluate()", "list[Finding]"],
        ["3. Score", "calculate_risk_score()", "float (0–100)"],
        ["4. Aggregate", "ScanEngine.scan()", "ScanResult"],
    ]
    add_table(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(7.5),
        table_data,
        col_widths=[Inches(1.5), Inches(3.5), Inches(2.5)],
    )

    # Design choices
    add_text_box(
        slide,
        Inches(8.8),
        Inches(2.2),
        Inches(4),
        Inches(0.5),
        "Key Design Choices",
        font_size=18,
        color=ACCENT_BLUE,
        bold=True,
    )
    add_accent_line(slide, Inches(8.8), Inches(2.75), Inches(2), ACCENT_BLUE)
    add_bullet_list(
        slide,
        Inches(8.8),
        Inches(3.0),
        Inches(4),
        Inches(3.5),
        [
            "Frozen dataclasses for\nimmutable scan results",
            "Registry pattern for O(1)\nrule lookup by resource type",
            "Dual scan paths\n(engine + legacy regex)",
            "Automatic deduplication\nof findings",
        ],
        font_size=14,
        color=LIGHT_GRAY,
        spacing=Pt(14),
    )


def slide_rules(prs):
    """Slide 6: Security Rules."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide,
        "Security Rules",
        "15 rules covering critical AWS infrastructure patterns",
    )

    # Engine rules table
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.0),
        Inches(5),
        Inches(0.4),
        "Engine Rules (5)",
        font_size=16,
        color=ACCENT_BLUE,
        bold=True,
    )
    table_data = [
        ["Rule", "Severity", "Catches"],
        ["S3 Public Access", "Critical", "Public ACL on S3 buckets"],
        ["Open SSH", "Critical", "Port 22 open to 0.0.0.0/0"],
        ["Public Database", "Critical", "RDS publicly_accessible = true"],
        ["Wildcard IAM", "Critical", 'Action: "*" in policies'],
        ["Disabled Encryption", "High", "storage_encrypted = false"],
    ]
    add_table(
        slide,
        Inches(0.8),
        Inches(2.5),
        Inches(7),
        table_data,
        col_widths=[Inches(2.2), Inches(1.3), Inches(3.5)],
    )

    # Seed rules
    add_text_box(
        slide,
        Inches(8.3),
        Inches(2.0),
        Inches(4.5),
        Inches(0.4),
        "Seed Guardrails (+10)",
        font_size=16,
        color=ACCENT_TEAL,
        bold=True,
    )
    seed_rules = [
        "Default VPC detection",
        "CloudWatch logging",
        "EBS volume encryption",
        "RDS backup retention",
        "Security group egress",
        "S3 versioning",
        "IAM user policies",
        "MFA enforcement",
        "VPC flow logs",
        "KMS key rotation",
    ]
    add_bullet_list(
        slide,
        Inches(8.3),
        Inches(2.5),
        Inches(4.5),
        Inches(4),
        seed_rules,
        font_size=13,
        color=LIGHT_GRAY,
        bullet_color=ACCENT_TEAL,
        spacing=Pt(4),
    )


def slide_scoring(prs):
    """Slide 7: Risk Scoring."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Risk Scoring", "Weighted severity model with 0–100 scale")

    # Scoring table
    table_data = [
        ["Severity", "Weight", "Rationale"],
        ["Critical", "10.0", "Immediate exploitation risk"],
        ["High", "7.0", "Significant security gap"],
        ["Medium", "4.0", "Best practice violation"],
        ["Low", "1.0", "Minor improvement"],
        ["Info", "0.5", "Advisory only"],
    ]
    add_table(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(5.5),
        table_data,
        col_widths=[Inches(1.5), Inches(1.2), Inches(2.8)],
    )

    # Formula
    add_text_box(
        slide,
        Inches(7),
        Inches(2.2),
        Inches(5.5),
        Inches(0.5),
        "Formula",
        font_size=18,
        color=ACCENT_BLUE,
        bold=True,
    )
    add_accent_line(slide, Inches(7), Inches(2.75), Inches(2), ACCENT_BLUE)

    formula_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(3.0), Inches(5.5), Inches(1)
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = DARK_BLUE
    formula_box.line.color.rgb = ACCENT_BLUE
    formula_box.line.width = Pt(1)
    add_text_box(
        slide,
        Inches(7.3),
        Inches(3.15),
        Inches(4.9),
        Inches(0.7),
        "score = (weighted_sum / max_possible) × 100",
        font_size=18,
        color=ACCENT_TEAL,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        font_name="Consolas",
    )

    # Score examples as cards
    add_text_box(
        slide,
        Inches(7),
        Inches(4.5),
        Inches(5.5),
        Inches(0.4),
        "Score Interpretation",
        font_size=16,
        color=WHITE,
        bold=True,
    )

    scores = [
        ("0", "Clean", GREEN),
        ("25", "Low Risk", YELLOW),
        ("60", "High Risk", ORANGE),
        ("100", "Critical", RED),
    ]
    for i, (val, label, color) in enumerate(scores):
        x = Inches(7 + i * 1.35)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.0), Inches(1.2), Inches(1.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BLUE
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_text_box(
            slide,
            x,
            Inches(5.1),
            Inches(1.2),
            Inches(0.6),
            val,
            font_size=24,
            color=color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x,
            Inches(5.7),
            Inches(1.2),
            Inches(0.4),
            label,
            font_size=11,
            color=MID_GRAY,
            alignment=PP_ALIGN.CENTER,
        )


def slide_frontend(prs):
    """Slide 8: Frontend Dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide, "Frontend Dashboard", "React 19 + TypeScript + Tailwind CSS + Recharts"
    )

    pages = [
        ("📊", "Dashboard", "Stats cards, severity bar chart,\nrecent scans overview"),
        ("📋", "Scan List", "Paginated table with\nstatus & file type filters"),
        ("📤", "New Scan", "Upload .tf file or\npaste content directly"),
        ("🔍", "Scan Detail", "Findings by severity,\nrisk score badge"),
        ("⚙️", "Guardrails", "Security rule\nmanagement interface"),
    ]

    for i, (icon, title, desc) in enumerate(pages):
        x = Inches(0.6 + i * 2.5)
        y = Inches(2.5)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(3.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BLUE
        card.line.color.rgb = RGBColor(0x33, 0x44, 0x55)
        card.line.width = Pt(1)

        add_text_box(
            slide,
            x,
            y + Inches(0.3),
            Inches(2.3),
            Inches(0.7),
            icon,
            font_size=36,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            y + Inches(1.1),
            Inches(2),
            Inches(0.4),
            title,
            font_size=16,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x + Inches(0.15),
            y + Inches(1.6),
            Inches(2),
            Inches(1.2),
            desc,
            font_size=12,
            color=MID_GRAY,
            alignment=PP_ALIGN.CENTER,
        )

    # Tech stack bar
    add_text_box(
        slide,
        Inches(0.8),
        Inches(6.2),
        Inches(11.5),
        Inches(0.4),
        "React 19  •  TypeScript  •  Tailwind CSS 3  •  Recharts  •  React Query  •  React Router 7",
        font_size=14,
        color=ACCENT_TEAL,
        alignment=PP_ALIGN.CENTER,
    )


def slide_tech_stack(prs):
    """Slide 9: Tech Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide, "Tech Stack", "Modern async Python + React with full DevOps pipeline"
    )

    table_data = [
        ["Layer", "Technology", "Purpose"],
        ["API Framework", "FastAPI", "Async REST API with OpenAPI docs"],
        ["ORM", "SQLAlchemy 2.x Async", "Async database operations"],
        ["Database", "SQLite + aiosqlite", "Zero-config persistence"],
        ["Validation", "Pydantic v2", "Request/response schemas"],
        ["Frontend", "React 19 + TypeScript", "Type-safe UI components"],
        ["Styling", "Tailwind CSS 3", "Utility-first responsive design"],
        ["Charts", "Recharts", "Severity distribution visualization"],
        ["Server State", "@tanstack/react-query", "Cache, sync, background refresh"],
        ["Testing", "pytest + Vitest + RTL", "Backend + frontend coverage"],
        ["CI/CD", "GitHub Actions", "5-job automated pipeline"],
        ["Containers", "Docker + Compose", "Multi-stage production builds"],
    ]
    add_table(
        slide,
        Inches(1.5),
        Inches(2.0),
        Inches(10.3),
        table_data,
        col_widths=[Inches(2), Inches(3.3), Inches(5)],
    )


def slide_quality(prs):
    """Slide 10: Testing & Quality."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide,
        "Testing & Quality Assurance",
        "Automated testing, linting, and security audit pipeline",
    )

    # Metric cards
    metrics = [
        ("Backend Tests", "97", ACCENT_BLUE),
        ("Coverage", "94.65%", GREEN),
        ("Frontend Tests", "24", ACCENT_TEAL),
        ("Lint Violations", "0", GREEN),
        ("Security Findings", "26 / 10 fixed", ORANGE),
    ]
    for i, (title, value, color) in enumerate(metrics):
        x = Inches(0.6 + i * 2.5)
        add_card(slide, x, Inches(2.2), Inches(2.2), Inches(1.3), title, value, color)

    # Testing approach
    add_text_box(
        slide,
        Inches(0.8),
        Inches(4.0),
        Inches(5),
        Inches(0.4),
        "Testing Approach",
        font_size=18,
        color=ACCENT_BLUE,
        bold=True,
    )
    add_accent_line(slide, Inches(0.8), Inches(4.45), Inches(2), ACCENT_BLUE)
    add_bullet_list(
        slide,
        Inches(0.8),
        Inches(4.6),
        Inches(5.5),
        Inches(2.5),
        [
            "Async integration tests with httpx.AsyncClient",
            "In-memory SQLite for test isolation",
            "React Testing Library for behavior testing",
            "Pre-commit hooks prevent regressions",
        ],
        font_size=14,
        color=LIGHT_GRAY,
    )

    # Quality gates
    add_text_box(
        slide,
        Inches(7),
        Inches(4.0),
        Inches(5),
        Inches(0.4),
        "Quality Gates",
        font_size=18,
        color=ACCENT_TEAL,
        bold=True,
    )
    add_accent_line(slide, Inches(7), Inches(4.45), Inches(2), ACCENT_TEAL)
    add_bullet_list(
        slide,
        Inches(7),
        Inches(4.6),
        Inches(5.5),
        Inches(2.5),
        [
            "ruff — linting (E, F, W, I, N, UP, S, B)",
            "black — code formatting",
            "mypy — static type checking",
            "tsc — TypeScript type checking",
        ],
        font_size=14,
        color=LIGHT_GRAY,
        bullet_color=ACCENT_TEAL,
    )


def slide_devops(prs):
    """Slide 11: DevOps & Deployment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide, "DevOps & Deployment", "Production-ready containerized deployment"
    )

    # CI Pipeline
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(5),
        Inches(0.4),
        "GitHub Actions CI — 5 Jobs",
        font_size=16,
        color=ACCENT_BLUE,
        bold=True,
    )

    ci_jobs = [
        "backend-lint",
        "backend-test",
        "frontend-lint",
        "frontend-test",
        "docker-build",
    ]
    for i, job in enumerate(ci_jobs):
        x = Inches(0.8 + i * 2.2)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.8), Inches(2), Inches(0.6)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BLUE
        box.line.color.rgb = GREEN
        box.line.width = Pt(1)
        add_text_box(
            slide,
            x,
            Inches(2.88),
            Inches(2),
            Inches(0.4),
            f"✓ {job}",
            font_size=12,
            color=GREEN,
            alignment=PP_ALIGN.CENTER,
        )

    # Docker
    add_text_box(
        slide,
        Inches(0.8),
        Inches(3.9),
        Inches(5.5),
        Inches(0.4),
        "Docker",
        font_size=16,
        color=ACCENT_TEAL,
        bold=True,
    )
    add_bullet_list(
        slide,
        Inches(0.8),
        Inches(4.4),
        Inches(5.5),
        Inches(2),
        [
            "Multi-stage builds (small images)",
            "Non-root container users",
            "Healthcheck endpoints",
            "docker-compose one-command deploy",
        ],
        font_size=14,
        color=LIGHT_GRAY,
        bullet_color=ACCENT_TEAL,
    )

    # Developer tools
    add_text_box(
        slide,
        Inches(7),
        Inches(3.9),
        Inches(5.5),
        Inches(0.4),
        "Developer Tools",
        font_size=16,
        color=ACCENT_BLUE,
        bold=True,
    )
    add_bullet_list(
        slide,
        Inches(7),
        Inches(4.4),
        Inches(5.5),
        Inches(2),
        [
            "Makefile with 12 shortcuts",
            "Pre-commit hooks (ruff, black, mypy)",
            "Secret detection in commits",
            ".env.example for easy setup",
        ],
        font_size=14,
        color=LIGHT_GRAY,
    )

    # Command examples
    cmd_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1)
    )
    cmd_box.fill.solid()
    cmd_box.fill.fore_color.rgb = RGBColor(0x0D, 0x11, 0x17)
    cmd_box.line.color.rgb = RGBColor(0x33, 0x44, 0x55)
    cmd_box.line.width = Pt(1)
    add_text_box(
        slide,
        Inches(1.2),
        Inches(6.15),
        Inches(10.5),
        Inches(0.7),
        "$ docker compose up --build -d    # Full stack in 60 seconds\n"
        "$ make test                       # Run all tests\n"
        "$ make lint                       # Run all linters",
        font_size=13,
        color=GREEN,
        font_name="Consolas",
    )


def slide_limitations(prs):
    """Slide 12: MVP Limitations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(
        slide, "MVP Tradeoffs & Limitations", "Pragmatic decisions for rapid delivery"
    )

    table_data = [
        ["Limitation", "Rationale", "Future Fix"],
        ["No authentication", "MVP scope — scanner logic first", "JWT / OAuth2"],
        ["SQLite only", "Zero-config, adequate for demo", "PostgreSQL option"],
        ["Regex-based parser", "No binary deps, pure Python", "HCL2 library"],
        ["Sync scanning", "Simple flow for single files", "Background tasks"],
        ["AWS rules only", "Most common cloud provider", "Azure / GCP rules"],
        ["No CloudFormation", "Terraform-first approach", "YAML / JSON parser"],
    ]
    add_table(
        slide,
        Inches(1),
        Inches(2.2),
        Inches(11.3),
        table_data,
        col_widths=[Inches(2.5), Inches(4.3), Inches(4.5)],
    )


def slide_roadmap(prs):
    """Slide 13: Future Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Future Roadmap", "Planned enhancements beyond MVP")

    # Near-term
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.2),
        Inches(5),
        Inches(0.4),
        "Near-term",
        font_size=18,
        color=ACCENT_BLUE,
        bold=True,
    )
    add_accent_line(slide, Inches(0.8), Inches(2.65), Inches(1.5), ACCENT_BLUE)
    add_bullet_list(
        slide,
        Inches(0.8),
        Inches(2.9),
        Inches(5.5),
        Inches(3),
        [
            "API key authentication",
            "CloudFormation YAML/JSON parsing",
            "Azure and GCP rule packs",
            "Export results as PDF/CSV",
            "Error boundary in React frontend",
        ],
        font_size=15,
        color=LIGHT_GRAY,
        spacing=Pt(10),
    )

    # Long-term
    add_text_box(
        slide,
        Inches(7),
        Inches(2.2),
        Inches(5),
        Inches(0.4),
        "Long-term",
        font_size=18,
        color=ACCENT_TEAL,
        bold=True,
    )
    add_accent_line(slide, Inches(7), Inches(2.65), Inches(1.5), ACCENT_TEAL)
    add_bullet_list(
        slide,
        Inches(7),
        Inches(2.9),
        Inches(5.5),
        Inches(3.5),
        [
            "JWT/OAuth2 with RBAC",
            "Async scan processing",
            "Multi-file / zip upload",
            "Scan trend comparison charts",
            "Webhook notifications",
            "PostgreSQL migration option",
            "Custom guardrail creation via UI",
        ],
        font_size=15,
        color=LIGHT_GRAY,
        bullet_color=ACCENT_TEAL,
        spacing=Pt(10),
    )


def slide_demo(prs):
    """Slide 14: Demo Flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_title(slide, "Demo Flow", "Live walkthrough of the scanning workflow")

    steps = [
        ("1", "Start Stack", "docker compose up\n--build -d"),
        ("2", "Open Dashboard", "http://localhost:3000\nView stats & charts"),
        ("3", "Upload .tf File", "New Scan → Upload\nvulnerable-infra.tf"),
        ("4", "View Findings", "Scan Detail → Findings\ngrouped by severity"),
        ("5", "Check Score", "Risk score badge\n0–100 weighted"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 2.5)
        y = Inches(2.8)

        # Step number
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(0.9), Inches(0.9)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()
        add_text_box(
            slide,
            x + Inches(0.65),
            y + Inches(0.12),
            Inches(0.9),
            Inches(0.65),
            num,
            font_size=28,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Connector
        if i < 4:
            add_text_box(
                slide,
                x + Inches(1.7),
                y + Inches(0.2),
                Inches(0.7),
                Inches(0.5),
                "→",
                font_size=24,
                color=MID_GRAY,
                alignment=PP_ALIGN.CENTER,
            )

        add_text_box(
            slide,
            x + Inches(0.1),
            y + Inches(1.1),
            Inches(2),
            Inches(0.4),
            title,
            font_size=16,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_text_box(
            slide,
            x + Inches(0.1),
            y + Inches(1.6),
            Inches(2),
            Inches(0.9),
            desc,
            font_size=12,
            color=MID_GRAY,
            alignment=PP_ALIGN.CENTER,
        )

    # Sample files note
    cmd_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.8)
    )
    cmd_box.fill.solid()
    cmd_box.fill.fore_color.rgb = DARK_BLUE
    cmd_box.line.color.rgb = ACCENT_TEAL
    cmd_box.line.width = Pt(1)
    add_text_box(
        slide,
        Inches(1.8),
        Inches(5.7),
        Inches(9.7),
        Inches(0.5),
        "📁  Sample files: samples/vulnerable-infra.tf (11 vulns)  •  "
        "samples/multi-service-vulnerable.tf (9 vulns)",
        font_size=14,
        color=ACCENT_TEAL,
        alignment=PP_ALIGN.CENTER,
    )


def slide_closing(prs):
    """Slide 15: Closing Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    # Top accent
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_text_box(
        slide,
        Inches(1.5),
        Inches(1.2),
        Inches(10.3),
        Inches(0.8),
        "Enterprise Security Guardrail Auditor",
        font_size=36,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_accent_line(slide, Inches(5.2), Inches(2.1), Inches(3), ACCENT_BLUE)

    checks = [
        "✅  Automated IaC security scanning — 15 rules",
        "✅  Weighted risk scoring — 0 to 100",
        "✅  Visual dashboard with severity charts",
        "✅  RESTful API with OpenAPI documentation",
        "✅  94.65% test coverage, zero lint violations",
        "✅  Production-ready Docker deployment",
        "✅  Extensible rule system for future growth",
    ]
    for i, item in enumerate(checks):
        add_text_box(
            slide,
            Inches(3),
            Inches(2.6 + i * 0.48),
            Inches(7.3),
            Inches(0.45),
            item,
            font_size=18,
            color=LIGHT_GRAY,
            alignment=PP_ALIGN.LEFT,
        )

    # Repo link
    link_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(6.0), Inches(6.3), Inches(0.7)
    )
    link_box.fill.solid()
    link_box.fill.fore_color.rgb = DARK_BLUE
    link_box.line.color.rgb = ACCENT_BLUE
    link_box.line.width = Pt(1)
    add_text_box(
        slide,
        Inches(3.5),
        Inches(6.1),
        Inches(6.3),
        Inches(0.5),
        "github.com/SimeonDee/Enterprise-Security-Guardrail-Auditor",
        font_size=16,
        color=ACCENT_BLUE,
        bold=False,
        alignment=PP_ALIGN.CENTER,
    )

    # Bottom accent
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def slide_thank_you(prs):
    """Slide 16: Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    # Top accent
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_text_box(
        slide,
        Inches(1.5),
        Inches(2.2),
        Inches(10.3),
        Inches(1.2),
        "Thank You",
        font_size=54,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_accent_line(slide, Inches(5.2), Inches(3.5), Inches(3), ACCENT_BLUE)

    add_text_box(
        slide,
        Inches(2),
        Inches(4.0),
        Inches(9.3),
        Inches(0.6),
        "Questions & Discussion",
        font_size=24,
        color=MID_GRAY,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Inches(2),
        Inches(5.2),
        Inches(9.3),
        Inches(0.5),
        "Adedoyin Simeon Adeyemi",
        font_size=18,
        color=LIGHT_GRAY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Inches(2),
        Inches(5.7),
        Inches(9.3),
        Inches(0.5),
        "github.com/SimeonDee/Enterprise-Security-Guardrail-Auditor",
        font_size=16,
        color=ACCENT_BLUE,
        alignment=PP_ALIGN.CENTER,
    )

    # Bottom accent
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)  # 1: Title
    slide_problem(prs)  # 2: Problem Statement
    slide_solution(prs)  # 3: Solution Overview
    slide_architecture(prs)  # 4: Architecture
    slide_scanner(prs)  # 5: Scanner Engine
    slide_rules(prs)  # 6: Security Rules
    slide_scoring(prs)  # 7: Risk Scoring
    slide_frontend(prs)  # 8: Frontend Dashboard
    slide_tech_stack(prs)  # 9: Tech Stack
    slide_quality(prs)  # 10: Testing & QA
    slide_devops(prs)  # 11: DevOps & Deployment
    slide_limitations(prs)  # 12: MVP Limitations
    slide_roadmap(prs)  # 13: Roadmap
    slide_demo(prs)  # 14: Demo Flow
    slide_closing(prs)  # 15: Closing Summary
    slide_thank_you(prs)  # 16: Thank You

    output_path = os.path.join(
        os.path.dirname(__file__), "..", "docs", "presentation.pptx"
    )
    output_path = os.path.abspath(output_path)
    prs.save(output_path)
    print(f"✓ Presentation saved to {output_path}")
    print(f"  {len(prs.slides)} slides generated")


if __name__ == "__main__":
    main()
